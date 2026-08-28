#!/usr/bin/env python3
"""Build a project-status PowerPoint from the UID QR Generator evidence."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "output" / "UID_QR_Generator_Development_Metrics.pptx"

NAVY = RGBColor(19, 48, 55)
TEAL = RGBColor(20, 107, 104)
MINT = RGBColor(221, 239, 229)
GOLD = RGBColor(241, 188, 81)
CORAL = RGBColor(218, 111, 87)
INK = RGBColor(30, 42, 44)
MUTED = RGBColor(101, 119, 121)
PALE = RGBColor(247, 249, 246)
WHITE = RGBColor(255, 255, 255)


def box(slide, x, y, w, h, fill=WHITE, line=None, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def text(slide, value, x, y, w, h, size=18, color=INK, bold=False, font="Aptos", align=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = value
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    if align:
        paragraph.alignment = align
    return shape


def title(slide, kicker, heading, number):
    text(slide, kicker.upper(), 0.65, 0.42, 4.5, 0.25, 9, TEAL, True, "Aptos")
    text(slide, heading, 0.65, 0.73, 10.8, 0.62, 27, NAVY, True, "Aptos Display")
    text(slide, f"UID QR GENERATOR  /  {number:02d}", 11.55, 0.48, 1.25, 0.25, 8, MUTED, True, "Aptos", PP_ALIGN.RIGHT)
    box(slide, 0.65, 1.48, 12.0, 0.02, TEAL)


def metric(slide, x, y, value, label, accent=TEAL, width=2.75):
    box(slide, x, y, width, 1.16, WHITE, RGBColor(224, 232, 228), True)
    box(slide, x, y + 1.12, width, 0.04, accent)
    text(slide, value, x + 0.16, y + 0.15, width - 0.32, 0.48, 25, NAVY, True, "Aptos Display")
    text(slide, label.upper(), x + 0.16, y + 0.72, width - 0.32, 0.22, 8, MUTED, True)


def bullets(slide, items, x, y, w, h, size=16, color=INK):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.06)
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(10)
        p.text = "•  " + p.text
    return shape


def new_slide(prs, bg=PALE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = bg
    return slide


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = new_slide(prs, NAVY)
    box(slide, 0.68, 0.78, 0.72, 0.72, TEAL, TEAL, True)
    text(slide, "U", 0.68, 0.79, 0.72, 0.7, 27, WHITE, True, "Aptos Display", PP_ALIGN.CENTER)
    text(slide, "UID QR Generator", 0.72, 1.82, 8.8, 0.78, 36, WHITE, True, "Aptos Display")
    text(slide, "Development metrics, architecture and effort invested", 0.75, 2.72, 8.8, 0.44, 19, MINT, False)
    box(slide, 0.75, 4.18, 11.75, 0.02, GOLD)
    text(slide, "A deterministic identity and QR workflow with an explainable operational panel", 0.75, 4.45, 8.7, 0.7, 22, WHITE, False)
    text(slide, "PROJECT REVIEW  ·  23 AUG 2026", 0.75, 6.72, 5, 0.25, 9, MINT, True)

    slide = new_slide(prs)
    title(slide, "Executive summary", "A complete path from source workbook to operational view", 2)
    metric(slide, 0.7, 1.95, "122", "synthetic assets processed", TEAL)
    metric(slide, 3.7, 1.95, "109", "serviceable assets", RGBColor(103, 179, 150))
    metric(slide, 6.7, 1.95, "13", "attention items", CORAL)
    metric(slide, 9.7, 1.95, "7", "formations represented", GOLD)
    bullets(slide, [
        "Built a repeatable UID assignment workflow around an external secret key.",
        "Added duplicate detection, historical mapping preservation and integrity checks.",
        "Produced distributable QR assets plus a data-driven local panel that mirrors Power BI interactions.",
        "Created a clean CSV handoff for Power BI slicers, measures and asset-level drill-through.",
    ], 0.9, 3.65, 11.3, 2.3, 17)

    slide = new_slide(prs)
    title(slide, "System delivered", "Five connected layers make the workflow usable", 3)
    layers = [
        ("01", "Input inspection", "Locate the ID field, identify the data row and read the workbook safely.", TEAL),
        ("02", "Validation gate", "Normalize IDs, stop on duplicates and prevent unsafe generation.", CORAL),
        ("03", "UID + QR generation", "Create deterministic HMAC-based UIDs and one QR PNG per UID.", GOLD),
        ("04", "Historical continuity", "Preserve previous mappings and append only genuinely new records.", RGBColor(103, 179, 150)),
        ("05", "Reporting surface", "Export Power BI CSV and drive the local interactive operations panel.", RGBColor(82, 136, 170)),
    ]
    for i, (num, head, body, color) in enumerate(layers):
        x = 0.78 + (i % 3) * 4.08
        y = 1.95 + (i // 3) * 2.25
        box(slide, x, y, 3.63, 1.72, WHITE, RGBColor(224, 232, 228), True)
        text(slide, num, x + 0.2, y + 0.18, 0.45, 0.32, 12, color, True)
        text(slide, head, x + 0.2, y + 0.58, 3.1, 0.34, 18, NAVY, True, "Aptos Display")
        text(slide, body, x + 0.2, y + 1.02, 3.18, 0.52, 11, MUTED)

    slide = new_slide(prs)
    title(slide, "Engineering metrics", "The current build is measurable and reproducible", 4)
    metric(slide, 0.7, 1.95, "32 B", "secret key length enforced", TEAL)
    metric(slide, 3.7, 1.95, "96-bit", "UID payload entropy", RGBColor(82, 136, 170))
    metric(slide, 6.7, 1.95, "100%", "input IDs mapped", RGBColor(103, 179, 150))
    metric(slide, 9.7, 1.95, "0", "duplicate UID collisions", CORAL)
    box(slide, 0.72, 3.62, 5.85, 2.36, MINT, MINT, True)
    text(slide, "VALIDATION COVERAGE", 1.0, 3.88, 3.8, 0.26, 9, TEAL, True)
    bullets(slide, ["Duplicate IDs checked before generation", "Historical UIDs compared after rerun", "Every generated UID checked for a QR file"], 1.0, 4.28, 5.0, 1.4, 14)
    box(slide, 6.85, 3.62, 5.75, 2.36, WHITE, RGBColor(224, 232, 228), True)
    text(slide, "DATASET SHAPE", 7.15, 3.88, 3.8, 0.26, 9, TEAL, True)
    text(slide, "21", 7.15, 4.25, 1.0, 0.48, 27, NAVY, True, "Aptos Display")
    text(slide, "source fields", 8.2, 4.36, 2.0, 0.24, 12, MUTED)
    text(slide, "6", 7.15, 5.02, 1.0, 0.48, 27, NAVY, True, "Aptos Display")
    text(slide, "types  ·  7 formations  ·  13 units", 8.2, 5.13, 3.8, 0.24, 12, MUTED)

    slide = new_slide(prs)
    title(slide, "Security and continuity", "The key design decision was preserving trust across reruns", 5)
    box(slide, 0.8, 1.95, 3.55, 3.8, NAVY, NAVY, True)
    text(slide, "HMAC-SHA256", 1.12, 2.35, 2.9, 0.42, 22, WHITE, True, "Aptos Display")
    text(slide, "secret key + normalized ID", 1.12, 2.95, 2.9, 0.3, 13, MINT)
    text(slide, "→", 1.12, 3.6, 0.5, 0.5, 28, GOLD, True)
    text(slide, "deterministic UID", 1.12, 4.25, 2.9, 0.42, 22, WHITE, True, "Aptos Display")
    text(slide, "same input and key, same result", 1.12, 4.86, 2.9, 0.3, 13, MINT)
    bullets(slide, [
        "Secret key is external to Excel and never printed.",
        "Historical register remains authoritative.",
        "New records append without changing old assignments.",
    ], 4.85, 2.18, 7.2, 2.2, 17)
    box(slide, 4.85, 4.65, 7.2, 1.1, MINT, MINT, True)
    text(slide, "Operational outcome", 5.15, 4.88, 1.8, 0.24, 10, TEAL, True)
    text(slide, "Repeatable identity, auditable changes and safe distribution boundaries.", 7.05, 4.83, 4.55, 0.42, 15, NAVY, True)

    slide = new_slide(prs)
    title(slide, "Interactive reporting", "The local panel now behaves like a Power BI prototype", 6)
    metric(slide, 0.7, 1.95, "122", "rows available to filter", TEAL)
    metric(slide, 3.7, 1.95, "21", "columns exported to CSV", RGBColor(82, 136, 170))
    metric(slide, 6.7, 1.95, "3", "interactive views", GOLD)
    metric(slide, 9.7, 1.95, "2", "refreshable outputs", CORAL)
    boxes = [("Overview", "Readiness, capability mix, serviceability, formations and cost", TEAL), ("Units & subunits", "Formation slicer with holdings and ready counts", GOLD), ("Fleet register", "Search across IDs, names, OEMs and units", CORAL)]
    for i, (head, body, color) in enumerate(boxes):
        x = 0.85 + i * 4.08
        box(slide, x, 3.7, 3.65, 1.65, WHITE, RGBColor(224, 232, 228), True)
        box(slide, x, 3.7, 0.08, 1.65, color)
        text(slide, head, x + 0.28, 4.03, 3.05, 0.34, 18, NAVY, True, "Aptos Display")
        text(slide, body, x + 0.28, 4.52, 3.0, 0.55, 12, MUTED)
    text(slide, "Power BI handoff: dashboard/fleet_register.csv", 0.88, 5.9, 6.5, 0.3, 13, TEAL, True)
    text(slide, "Local simulation: dashboard/index.html", 7.2, 5.9, 5.0, 0.3, 13, TEAL, True)

    slide = new_slide(prs)
    title(slide, "Effort invested", "Engineering effort was concentrated in correctness and usability", 7)
    text(slide, "ESTIMATED EFFORT BREAKDOWN", 0.85, 1.92, 4.0, 0.25, 9, TEAL, True)
    efforts = [("Workflow design & notebook", 25, TEAL), ("Validation & continuity", 25, CORAL), ("QR and workbook outputs", 15, GOLD), ("Dashboard interaction model", 20, RGBColor(82, 136, 170)), ("Power BI data preparation", 15, RGBColor(103, 179, 150))]
    y = 2.42
    for label, pct, color in efforts:
        text(slide, label, 0.9, y, 2.8, 0.3, 13, INK, True)
        box(slide, 3.7, y + 0.06, 6.3, 0.23, RGBColor(232, 237, 233))
        box(slide, 3.7, y + 0.06, 6.3 * pct / 100, 0.23, color)
        text(slide, f"{pct}%", 10.25, y, 0.55, 0.3, 13, NAVY, True, "Aptos", PP_ALIGN.RIGHT)
        y += 0.62
    box(slide, 0.88, 5.72, 11.5, 0.65, MINT, MINT, True)
    text(slide, "Estimate basis: relative allocation of development attention from the delivered modules; not a timesheet or billing record.", 1.12, 5.88, 11.0, 0.28, 11, MUTED)

    slide = new_slide(prs)
    title(slide, "Readiness and next steps", "The foundation is ready for a controlled Power BI rollout", 8)
    box(slide, 0.8, 1.95, 5.65, 3.75, MINT, MINT, True)
    text(slide, "DELIVERED", 1.12, 2.28, 2.0, 0.25, 9, TEAL, True)
    bullets(slide, ["Deterministic UID generation", "Duplicate and collision safeguards", "Historical mapping preservation", "QR image and Excel distribution", "Data-driven interactive panel", "Power BI-ready CSV export"], 1.1, 2.72, 4.8, 2.4, 15)
    box(slide, 6.85, 1.95, 5.65, 3.75, WHITE, RGBColor(224, 232, 228), True)
    text(slide, "RECOMMENDED NEXT", 7.17, 2.28, 2.7, 0.25, 9, TEAL, True)
    bullets(slide, ["Connect Power BI to the controlled export location", "Add role-based access and refresh ownership", "Replace synthetic inventory with approved production input", "Add deployment runbook and operational acceptance test"], 7.15, 2.72, 4.8, 2.4, 15)
    text(slide, "Success measure: trusted UID continuity plus fast operational filtering.", 0.88, 6.38, 11.4, 0.34, 18, NAVY, True, "Aptos Display", PP_ALIGN.CENTER)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote presentation: {OUT}")


if __name__ == "__main__":
    build()